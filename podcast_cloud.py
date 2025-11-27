# run command:
#                             streamlit run podcast_cloud.py
#                  directory setup: cd C:\users\oakhtar\documents\pyprojs_local
#   OPTIMIZED for deployment -
#
import streamlit as st
import os
import tempfile
import json
import requests
import shutil
import re
from pathlib import Path
from pydub import AudioSegment
from pydub.effects import high_pass_filter, low_pass_filter
from datetime import datetime

# --- TEXT PROCESSING ---
import PyPDF2
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
import yt_dlp

# --- AI CLIENT ---
from openai import OpenAI

# ================= CONFIGURATION =================
st.set_page_config(
    page_title="PodcastLM Studio", 
    page_icon="🎧", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# ================= SESSION STATE =================
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False
if "script_data" not in st.session_state:
    st.session_state.script_data = None
if "source_text" not in st.session_state:
    st.session_state.source_text = ""
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "notebook_content" not in st.session_state:
    st.session_state.notebook_content = f"# 📓 Research Notebook\n**Session Started:** {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n"

# ================= AUTHENTICATION =================
def check_password():
    user_pass = st.session_state.get("password_input", "")
    correct_pass = st.secrets.get("APP_PASSWORD")
    if correct_pass and user_pass == correct_pass:
        st.session_state.authenticated = True
    else:
        st.error("❌ Incorrect Password")

if not st.session_state.authenticated:
    st.title("🔒 Studio Login")
    st.text_input("Enter Password", type="password", key="password_input", on_change=check_password)
    st.stop()

# ================= UTILS =================

def download_file_with_headers(url, save_path):
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
        response = requests.get(url, headers=headers, stream=True, timeout=15)
        if response.status_code == 200:
            with open(save_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)
            return True
        return False
    except: return False

def scrape_website(url):
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
        response = requests.get(url, headers=headers, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')
        for script in soup(["script", "style", "header", "footer", "nav"]):
            script.decompose()
        return soup.get_text()
    except: return None

def extract_text_from_files(files, client=None):
    text = ""
    for file in files:
        try:
            name = file.name.lower()
            if name.endswith(".pdf"):
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages: text += page.extract_text() + "\n"
            elif name.endswith(".docx"):
                doc = docx.Document(file)
                for para in doc.paragraphs: text += para.text + "\n"
            elif name.endswith(".pptx"):
                prs = Presentation(file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"): text += shape.text + "\n"
            elif name.endswith(".txt"):
                text = file.getvalue().decode("utf-8")
            elif name.endswith((".mp3", ".mp4", ".wav", ".m4a", ".mpeg", ".webm")):
                if client:
                    with st.spinner(f"Transcribing {name}..."):
                        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(name)[1]) as tmp_file:
                            tmp_file.write(file.getvalue())
                            tmp_path = tmp_file.name
                        with open(tmp_path, "rb") as audio_file:
                            transcript = client.audio.transcriptions.create(model="whisper-1", file=audio_file)
                        text += transcript.text + "\n"
                        os.remove(tmp_path)
                else:
                    st.warning(f"Skipped {file.name}: API Key required.")
            else:
                st.warning(f"Skipped {file.name}: Unsupported format.")
        except Exception as e:
            st.error(f"Error reading {file.name}: {e}")
    return text

def download_and_transcribe_video(url, client):
    try:
        with tempfile.TemporaryDirectory() as tmp_dir:
            ydl_opts = {
                'format': 'bestaudio/best',
                'outtmpl': os.path.join(tmp_dir, 'audio.%(ext)s'),
                'postprocessors': [{'key': 'FFmpegExtractAudio','preferredcodec': 'mp3','preferredquality': '128'}],
                'quiet': True, 'no_warnings': True, 'nocheckcertificate': True,
                'http_headers': {'User-Agent': 'Mozilla/5.0'}
            }
            with yt_dlp.YoutubeDL(ydl_opts) as ydl: ydl.download([url])
            audio_path = os.path.join(tmp_dir, "audio.mp3")
            if not os.path.exists(audio_path): return None, "Download failed."
            if os.path.getsize(audio_path) / (1024*1024) > 24: return None, "Video too long (>25MB)."
            with open(audio_path, "rb") as f:
                transcript = client.audio.transcriptions.create(model="whisper-1", file=f)
            return transcript.text, None
    except Exception as e: return None, str(e)

def generate_audio_openai(client, text, voice, filename, speed=1.0):
    try:
        response = client.audio.speech.create(model="tts-1", voice=voice, input=text, speed=speed)
        response.stream_to_file(filename)
        return True
    except: return False

# ================= SIDEBAR =================

with st.sidebar:
    st.title("🎛️ Studio Settings")
    api_key = st.secrets.get("OPENAI_API_KEY") or st.text_input("OpenAI API Key", type="password")
    
    privacy_mode = st.toggle("🛡️ Privacy Mode", value=False)
    
    if st.button("🗑️ New Session (Clear All)"):
        st.session_state.chat_history = []
        st.session_state.notebook_content = f"# 📓 Research Notebook\n**Session Started:** {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n"
        st.session_state.source_text = ""
        st.session_state.script_data = None
        st.rerun()
        
    st.divider()
    
    st.subheader("🌍 Localization")
    language = st.selectbox("Output Language", [
        "English (US)", "English (UK)", "Spanish (Spain)", "Spanish (LatAm)", 
        "French", "German", "Italian", "Portuguese", "Portuguese (Brazil)",
        "Japanese", "Chinese (Mandarin)", "Korean", "Hindi", "Urdu", "Arabic", "Russian",
        "Turkish", "Dutch", "Polish", "Swedish", "Danish", "Norwegian", "Finnish",
        "Greek", "Czech", "Romanian", "Indonesian", "Vietnamese", "Thai", "Hebrew"
    ])
    
    length_option = st.select_slider("Duration", ["Short (2 min)", "Medium (5 min)", "Long (15 min)", "Extra Long (30 min)"])

    st.subheader("🎭 Hosts")
    host1_persona = st.text_input("Host 1 Persona", "Male, curious, slightly skeptical")
    host2_persona = st.text_input("Host 2 Persona", "Female, enthusiastic expert, fast talker")
    
    voice_style = st.selectbox("Voice Pair", ["Dynamic (Alloy & Nova)", "Calm (Onyx & Shimmer)", "Formal (Echo & Fable)"])
    voice_map = {"Dynamic (Alloy & Nova)": ("alloy", "nova"), "Calm (Onyx & Shimmer)": ("onyx", "shimmer"), "Formal (Echo & Fable)": ("echo", "fable")}

    st.divider()
    st.subheader("🎵 Music & Branding")
    
    bg_source = st.radio("Background Music", ["Presets", "Upload Custom", "None"], horizontal=True)
    music_ramp_up = st.checkbox("🎵 Start Music 5s Before Dialogue", value=False)

    selected_bg_url = None
    uploaded_bg_file = None
    if bg_source == "Presets":
        music_choice = st.selectbox("Track", ["Lo-Fi (Study)", "Upbeat (Morning)", "Ambient (News)", "Cinematic (Deep)"])
        music_urls = {
            "Lo-Fi (Study)": "https://cdn.pixabay.com/download/audio/2022/05/27/audio_1808fbf07a.mp3?filename=lofi-study-112191.mp3",
            "Upbeat (Morning)": "https://cdn.pixabay.com/download/audio/2024/05/24/audio_95e3f5f471.mp3?filename=good-morning-206098.mp3",
            "Ambient (News)": "https://cdn.pixabay.com/download/audio/2022/03/10/audio_c8c8a73467.mp3?filename=ambient-piano-10226.mp3",
            "Cinematic (Deep)": "https://cdn.pixabay.com/download/audio/2022/03/22/audio_c2b86c77ce.mp3?filename=cinematic-atmosphere-score-2-21266.mp3"
        }
        selected_bg_url = music_urls[music_choice]
    elif bg_source == "Upload Custom":
        uploaded_bg_file = st.file_uploader("Upload Loop (MP3/WAV)", type=["mp3", "wav"])

    with st.expander("Intro/Outro Clips"):
        uploaded_intro = st.file_uploader("Intro (Plays Once)", type=["mp3", "wav"])
        uploaded_outro = st.file_uploader("Outro (Plays Once)", type=["mp3", "wav"])

# ================= MAIN APP =================
st.title("🎧 PodcastLM Studio")

tab1, tab2, tab3, tab4 = st.tabs(["1. Source Material", "2. 🤖 AI Research Assistant", "3. Script Editor", "4. Audio Production"])

# --- TAB 1: INPUT ---
with tab1:
    st.info("Upload content here. This drives both the **Podcast** and the **Chatbot**.")
    input_type = st.radio("Input Type", ["📂 Files", "🔗 Web URL", "📺 Video URL", "📝 Text"], horizontal=True)
    new_text = ""
    
    if input_type == "📂 Files":
        files = st.file_uploader("Upload", accept_multiple_files=True)
        if files and st.button("Process Files"):
            with st.spinner("Processing uploaded files..."):
                client = OpenAI(api_key=api_key) if api_key else None
                new_text = extract_text_from_files(files, client)
            
    elif input_type == "🔗 Web URL":
        url = st.text_input("Enter Article URL")
        if url and st.button("Scrape Website"): 
            with st.spinner("Scraping..."):
                scraped = scrape_website(url)
                if scraped: new_text = scraped
                else: st.error("Blocked by website.")
                
    elif input_type == "📺 Video URL":
        vid_url = st.text_input("Enter Video URL")
        if vid_url and st.button("Transcribe"):
            if api_key:
                client = OpenAI(api_key=api_key)
                with st.spinner("Downloading and Transcribing Video..."):
                    text, err = download_and_transcribe_video(vid_url, client)
                    if text: new_text = text
                    else: st.error(err)
            else: st.error("API Key Required")
    
    elif input_type == "📝 Text":
        new_text = st.text_area("Paste Text", height=300)

    if new_text and new_text != st.session_state.source_text:
        st.session_state.source_text = new_text
        st.session_state.chat_history = [] 
        timestamp = datetime.now().strftime("%H:%M:%S")
        st.session_state.notebook_content += f"\n---\n### 📥 New Source Loaded ({timestamp})\n*Source Type: {input_type}*\n\n"
        st.success("✅ Source text loaded!")

    if st.session_state.source_text:
        with st.expander("View Source Text"):
            st.text_area("Content", st.session_state.source_text, height=150, disabled=True)

# --- TAB 2: CHAT ---
with tab2:
    col_chat, col_notes = st.columns([1, 1])
    with col_chat:
        st.subheader("💬 Active Chat")
        if not st.session_state.source_text: st.warning("Load source text first.")
        else:
            for message in st.session_state.chat_history:
                with st.chat_message(message["role"]): st.markdown(message["content"])
            if prompt := st.chat_input("Ask a question..."):
                if api_key:
                    st.session_state.chat_history.append({"role": "user", "content": prompt})
                    st.session_state.notebook_content += f"**Q:** {prompt}\n\n"
                    with st.chat_message("user"): st.markdown(prompt)
                    with st.chat_message("assistant"):
                        client = OpenAI(api_key=api_key)
                        stream = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=[
                                {"role": "system", "content": "Answer based ONLY on source text."},
                                {"role": "user", "content": f"Source: {st.session_state.source_text[:30000]}"},
                                {"role": "user", "content": prompt}
                            ], stream=True)
                        response = st.write_stream(stream)
                    st.session_state.chat_history.append({"role": "assistant", "content": response})
                    st.session_state.notebook_content += f"**A:** {response}\n\n"
                    st.rerun()

    with col_notes:
        st.subheader("📓 Research Notebook")
        st.caption("Auto-saves Q&A. Editable.")
        updated_notebook = st.text_area("Notebook Content", value=st.session_state.notebook_content, height=600, key="notebook_area")
        if updated_notebook != st.session_state.notebook_content: st.session_state.notebook_content = updated_notebook
        st.download_button("💾 Save Notebook (.md)", st.session_state.notebook_content, f"notebook_{datetime.now().strftime('%Y%m%d_%H%M')}.md")

# --- TAB 3: SCRIPT ---
with tab3:
    st.markdown("### 🎬 Director Mode")
    col_dir, col_call = st.columns([1, 1])
    
    with col_dir:
        user_instructions = st.text_area(
            "📢 Custom Instructions", 
            placeholder="e.g., 'Focus on the financial aspects', 'Make it funny'"
        )
        
    with col_call:
        st.markdown("#### 📞 Call-in Segment")
        caller_prompt = st.text_area(
            "Listener Question/Comment",
            placeholder="Type a question here. E.g., 'Hey, I disagree with your point about X...'"
        )
        st.caption("If filled, a 'Caller' will interrupt the show with this question.")

    if st.button("Generate Podcast Script", type="primary"):
        if not api_key or not st.session_state.source_text: st.error("Missing Key or Text")
        else:
            try:
                client = OpenAI(api_key=api_key)
                length_instr = "12-15 exchanges"
                if "Medium" in length_option: length_instr = "30 exchanges. Deep dive."
                elif "Long" in length_option: length_instr = "50 exchanges. Very detailed."
                elif "Extra Long" in length_option: length_instr = "80 exchanges. Comprehensive."

                call_in_instr = ""
                if caller_prompt:
                    call_in_instr = f"""
                    MANDATORY: Halfway through the script, include a 'Listener Call-in' segment.
                    The speaker label must be "Caller".
                    The Caller says: "{caller_prompt}".
                    Host 1 and Host 2 must react to this call and debate the caller's point.
                    """

                prompt = f"""
                Create a podcast script.
                Language: {language}
                Length: {length_instr}
                Host 1: {host1_persona}
                Host 2: {host2_persona}
                
                DIRECTOR NOTES: {user_instructions}
                {call_in_instr}
                
                Format: JSON {{ "title": "...", "dialogue": [ {{"speaker": "Host 1", "text": "..."}}, {{"speaker": "Caller", "text": "..."}} ] }}
                Text: {st.session_state.source_text[:35000]}
                """
                
                with st.spinner("Drafting Script..."):
                    res = client.chat.completions.create(
                        model="gpt-4o-mini", messages=[{"role": "user", "content": prompt}], response_format={"type": "json_object"}
                    )
                    st.session_state.script_data = json.loads(res.choices[0].message.content)
                    st.success("Ready!")
                    if privacy_mode: st.session_state.source_text = ""
            except Exception as e: st.error(f"Error: {e}")

    if st.session_state.script_data:
        data = st.session_state.script_data
        st.subheader(data.get('title', 'Podcast'))
        with st.form("edit"):
            new_d = []
            for i, l in enumerate(data['dialogue']):
                c1, c2 = st.columns([1, 5])
                # Add "Caller" to the dropdown options if it appears
                roles = ["Host 1", "Host 2"]
                if l['speaker'] == "Caller": roles.append("Caller")
                
                idx = 0
                if l['speaker'] == "Host 2": idx = 1
                elif l['speaker'] == "Caller": idx = 2
                
                spk = c1.selectbox("Role", roles, index=idx if idx < len(roles) else 0, key=f"s{i}")
                txt = c2.text_area("Line", l['text'], height=70, key=f"t{i}")
                new_d.append({"speaker": spk, "text": txt})
            if st.form_submit_button("Save"):
                st.session_state.script_data['dialogue'] = new_d
                st.success("Saved")

# --- TAB 4: AUDIO ---
with tab4:
    if st.session_state.script_data and st.button("🎙️ Start Production", type="primary"):
        if not api_key: st.stop()
        
        progress = st.progress(0)
        status = st.empty()
        client = OpenAI(api_key=api_key)
        m_voice, f_voice = voice_map[voice_style]
        
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            segs = []
            script = st.session_state.script_data['dialogue']
            
            for i, line in enumerate(script):
                status.text(f"Recording {i+1}/{len(script)}...")
                
                # Voice Selection Logic
                voice = m_voice
                if line['speaker'] == "Host 2": voice = f_voice
                elif line['speaker'] == "Caller": voice = "fable" # Distinct voice for caller
                
                f_path = str(tmp_path / f"line_{i}.mp3")
                if generate_audio_openai(client, line['text'], voice, f_path):
                    seg = AudioSegment.from_file(f_path)
                    
                    # Apply Telephone FX if Caller
                    if line['speaker'] == "Caller":
                        seg = high_pass_filter(seg, 300)
                        seg = low_pass_filter(seg, 3000)
                        
                    segs.append(seg)
                    segs.append(AudioSegment.silent(duration=350))
                progress.progress((i+1)/len(script))
            
            if segs:
                status.text("Mixing...")
                final = sum(segs)
                
                # Music Logic
                bg_seg = None
                try:
                    if bg_source == "Presets" and selected_bg_url:
                        if download_file_with_headers(selected_bg_url, tmp_path/"bg.mp3"): bg_seg = AudioSegment.from_file(tmp_path/"bg.mp3")
                    elif bg_source == "Upload Custom" and uploaded_bg_file:
                        with open(tmp_path/"bg.mp3", "wb") as f: f.write(uploaded_bg_file.getvalue())
                        bg_seg = AudioSegment.from_file(tmp_path/"bg.mp3")
                    
                    if bg_seg:
                        bg_seg = bg_seg - 22
                        if music_ramp_up: final = AudioSegment.silent(duration=5000) + final
                        while len(bg_seg) < len(final) + 5000: bg_seg += bg_seg
                        bg_seg = bg_seg[:len(final)+2000].fade_out(3000)
                        final = bg_seg.overlay(final)
                except: pass

                try:
                    if uploaded_intro:
                        with open(tmp_path/"in.mp3","wb") as f: f.write(uploaded_intro.getvalue())
                        final = AudioSegment.from_file(tmp_path/"in.mp3") + final
                    if uploaded_outro:
                        with open(tmp_path/"out.mp3","wb") as f: f.write(uploaded_outro.getvalue())
                        final = final + AudioSegment.from_file(tmp_path/"out.mp3")
                except: pass

                final.export(tmp_path/"master.mp3", format="mp3", bitrate="192k")
                with open(tmp_path/"master.mp3", "rb") as f: ab = f.read()
                
                status.success("Done!")
                st.audio(ab, format="audio/mp3")
                st.download_button("Download MP3", ab, "podcast.mp3", "audio/mp3")




